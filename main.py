#!/usr/bin/env python3
"""
AI-powered presentation video generator from PDF slides
"""

import os
import sys
from pathlib import Path
from typing import List, Optional
import logging

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

from dotenv import load_dotenv
from pdf2image import convert_from_path
from PIL import Image
from google import genai
from google.genai import types
import wave
import cv2
import numpy as np
try:
    from moviepy.editor import VideoFileClip, concatenate_videoclips
    from moviepy.audio.io.AudioFileClip import AudioFileClip
except ImportError:
    VideoFileClip = None
    concatenate_videoclips = None
    AudioFileClip = None
import tempfile
import shutil
from pptx import Presentation

# Load environment variables
load_dotenv()

class PresentationVideoGenerator:
    def __init__(self, project_name: str):
        self.genai_client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))
        self.project_name = project_name
        
        # プロジェクトごとのディレクトリ構造を作成
        self.base_output_dir = Path("output") / project_name
        self.txt_dir = self.base_output_dir / "txt"
        self.wav_dir = self.base_output_dir / "wav"
        self.mp4_dir = self.base_output_dir / "mp4"
        
        # ディレクトリを作成
        self.base_output_dir.mkdir(parents=True, exist_ok=True)
        self.txt_dir.mkdir(exist_ok=True)
        self.wav_dir.mkdir(exist_ok=True)
        self.mp4_dir.mkdir(exist_ok=True)
        
    def get_corresponding_pdf(self, pptx_path: str) -> str:
        """PPTXファイルに対応するPDFファイルのパスを取得"""
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = Path("slides/pdf") / f"{base_name}.pdf"
        
        if not pdf_path.exists():
            raise FileNotFoundError(f"Corresponding PDF file not found: {pdf_path}")
        
        logger.info(f"Found corresponding PDF: {pdf_path}")
        return str(pdf_path)
    
    def extract_speaker_notes_from_pptx(self, pptx_path: str) -> List[Optional[str]]:
        """PPTXファイルから発表者ノートを抽出"""
        logger.info(f"Extracting speaker notes from PPTX: {pptx_path}")
        try:
            presentation = Presentation(pptx_path)
            speaker_notes = []
            
            for i, slide in enumerate(presentation.slides, 1):
                note_text = None
                
                # 発表者ノートが存在するかチェック
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                        note_text = notes_slide.notes_text_frame.text.strip()
                        # 空文字列や空白のみの場合はNoneとして扱う
                        if not note_text:
                            note_text = None
                        else:
                            logger.info(f"Found speaker notes for slide {i}: {len(note_text)} characters")
                
                if note_text is None:
                    logger.info(f"No speaker notes found for slide {i}")
                
                speaker_notes.append(note_text)
            
            logger.info(f"Extracted speaker notes for {len(speaker_notes)} slides")
            return speaker_notes
            
        except Exception as e:
            logger.error(f"Error extracting speaker notes from PPTX: {e}")
            raise
    
    def extract_images_from_pdf(self, pdf_path: str) -> List[Image.Image]:
        """PDFから各ページの画像を抽出"""
        logger.info(f"Extracting images from PDF: {pdf_path}")
        try:
            images = convert_from_path(pdf_path, dpi=300)
            logger.info(f"Extracted {len(images)} pages from PDF")
            return images
        except Exception as e:
            logger.error(f"Error extracting images from PDF: {e}")
            raise
    
    def get_script_for_slide(self, speaker_note: Optional[str] = None) -> Optional[str]:
        """スライドの原稿を取得（発表者ノートのみ使用）"""
        
        # 発表者ノートがある場合はそれを使用
        if speaker_note is not None and speaker_note.strip():
            return speaker_note.strip()
        
        # 発表者ノートがない場合はNoneを返す
        return None
    
    
    def text_to_speech(self, text: str, output_path: str) -> str:
        """GeminiのTTSを使用してテキストを音声に変換"""
        logger.info(f"Converting text to speech: {len(text)} characters")
        
        try:
            # Gemini TTSの実装
            response = self.genai_client.models.generate_content(
                model="gemini-2.5-flash-preview-tts",
                contents=text,
                config=types.GenerateContentConfig(
                    response_modalities=["AUDIO"],
                    speech_config=types.SpeechConfig(
                        voice_config=types.VoiceConfig(
                            prebuilt_voice_config=types.PrebuiltVoiceConfig(
                                voice_name="Kore",
                            )
                        )
                    ),
                )
            )
            
            # 音声データを取得
            audio_data = response.candidates[0].content.parts[0].inline_data.data
            
            # WAVファイルとして保存
            with wave.open(output_path, "wb") as wf:
                wf.setnchannels(1)
                wf.setsampwidth(2)
                wf.setframerate(24000)
                wf.writeframes(audio_data)
            
            logger.info(f"Audio saved to: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Error in text-to-speech conversion: {e}")
            raise
    
    def create_video_from_slide_and_audio(self, image: Image.Image, audio_path: Optional[str], output_path: str, duration: float = 3.0) -> str:
        """スライド画像と音声（または無音）を合成して動画を作成"""
        if audio_path:
            logger.info(f"Creating video from slide and audio: {output_path}")
        else:
            logger.info(f"Creating silent video from slide: {output_path} (duration: {duration}s)")
        
        try:
            # 画像を一時ファイルに保存（サイズ調整）
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                # 高さを2で割り切れるように調整
                width, height = image.size
                if height % 2 != 0:
                    height -= 1
                if width % 2 != 0:
                    width -= 1
                
                # 画像をリサイズ
                resized_image = image.resize((width, height), Image.LANCZOS)
                resized_image.save(temp_file.name)
                temp_image_path = temp_file.name
            
            # FFmpegを使用して動画を作成
            import subprocess
            
            if audio_path:
                # 音声の長さを取得（ffprobeを使用）
                duration_cmd = [
                    'ffprobe', '-v', 'error', '-show_entries', 'format=duration', 
                    '-of', 'default=noprint_wrappers=1:nokey=1', audio_path
                ]
                duration_result = subprocess.run(duration_cmd, capture_output=True, text=True)
                duration = float(duration_result.stdout.strip())
                
                # FFmpegで画像と音声を合成
                ffmpeg_cmd = [
                    'ffmpeg', '-y',  # -y for overwrite
                    '-loop', '1',    # Loop the image
                    '-i', temp_image_path,  # Input image
                    '-i', audio_path,       # Input audio
                    '-c:v', 'libx264',      # Video codec
                    '-c:a', 'aac',          # Audio codec
                    '-t', str(duration),    # Duration
                    '-pix_fmt', 'yuv420p',  # Pixel format for compatibility
                    '-r', '1',              # Frame rate (1fps for static image)
                    '-shortest',            # Stop encoding when the shortest input ends
                    output_path
                ]
            else:
                # 無音動画を作成
                ffmpeg_cmd = [
                    'ffmpeg', '-y',  # -y for overwrite
                    '-loop', '1',    # Loop the image
                    '-i', temp_image_path,  # Input image
                    '-f', 'lavfi',   # Generate silent audio
                    '-i', 'anullsrc=channel_layout=mono:sample_rate=48000',  # Silent audio
                    '-c:v', 'libx264',      # Video codec
                    '-c:a', 'aac',          # Audio codec
                    '-t', str(duration),    # Duration
                    '-pix_fmt', 'yuv420p',  # Pixel format for compatibility
                    '-r', '1',              # Frame rate (1fps for static image)
                    '-shortest',            # Stop encoding when the shortest input ends
                    output_path
                ]
            
            result = subprocess.run(ffmpeg_cmd, capture_output=True, text=True)
            if result.returncode != 0:
                logger.error(f"FFmpeg error: {result.stderr}")
                raise subprocess.CalledProcessError(result.returncode, ffmpeg_cmd, result.stderr)
            
            # クリーンアップ
            os.unlink(temp_image_path)
            
            logger.info(f"Video created: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Error creating video: {e}")
            raise
    
    def combine_videos(self, video_paths: List[str], output_path: str) -> str:
        """複数の動画を結合して最終的なプレゼン動画を作成"""
        logger.info(f"Combining {len(video_paths)} videos into final presentation")
        
        try:
            # FFmpegを使用して動画を結合
            import subprocess
            
            # 入力リストファイルを作成
            input_list_path = self.base_output_dir / "input_list.txt"
            with open(input_list_path, 'w') as f:
                for video_path in video_paths:
                    f.write(f"file '{os.path.abspath(video_path)}'\n")
            
            # FFmpegで動画を結合
            ffmpeg_cmd = [
                'ffmpeg', '-y',  # -y for overwrite
                '-f', 'concat',  # Concat format
                '-safe', '0',    # Allow absolute paths
                '-i', str(input_list_path),  # Input list
                '-c', 'copy',    # Copy streams without re-encoding
                output_path
            ]
            
            subprocess.run(ffmpeg_cmd, check=True, capture_output=True)
            
            # クリーンアップ
            os.unlink(input_list_path)
            
            logger.info(f"Final presentation video created: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Error combining videos: {e}")
            raise
    
    def generate_presentation_video(self, file_path: str) -> str:
        """PPTXからプレゼン動画を生成するメイン処理"""
        logger.info(f"Starting presentation video generation from: {file_path}")
        
        try:
            # PPTX処理
            logger.info("Processing PPTX file")
            
            # 1. 発表者ノートを抽出
            speaker_notes = self.extract_speaker_notes_from_pptx(file_path)
            
            # 2. 対応するPDFファイルを検索
            corresponding_pdf_path = self.get_corresponding_pdf(file_path)
            
            # 3. PDFから画像を抽出
            images = self.extract_images_from_pdf(corresponding_pdf_path)
            
            video_paths = []
            
            # 4. 各スライドに対して処理
            for i, image in enumerate(images, 1):
                logger.info(f"Processing slide {i}/{len(images)}")
                
                # 5. 原稿取得（発表者ノートのみ使用）
                speaker_note = speaker_notes[i-1] if i <= len(speaker_notes) else None
                script = self.get_script_for_slide(speaker_note)
                
                if script:
                    # 原稿がある場合の処理
                    logger.info(f"Using speaker notes for slide {i}: {len(script)} characters")
                    
                    # 原稿をファイルに保存
                    script_path = self.txt_dir / f"slide_{i:03d}_script.txt"
                    with open(script_path, 'w', encoding='utf-8') as f:
                        f.write(script)
                    logger.info(f"Script saved to: {script_path}")
                    
                    # 6. 音声生成
                    audio_path = self.wav_dir / f"slide_{i:03d}_audio.wav"
                    self.text_to_speech(script, str(audio_path))
                    
                    # 7. 動画生成（音声付き）
                    video_path = self.mp4_dir / f"slide_{i:03d}_video.mp4"
                    self.create_video_from_slide_and_audio(image, str(audio_path), str(video_path))
                    video_paths.append(str(video_path))
                else:
                    # 原稿がない場合は3秒間の無音動画を作成
                    logger.info(f"No speaker notes for slide {i}, creating 3-second silent video")
                    
                    # 7. 無音動画生成（3秒）
                    video_path = self.mp4_dir / f"slide_{i:03d}_video.mp4"
                    self.create_video_from_slide_and_audio(image, None, str(video_path), duration=3.0)
                    video_paths.append(str(video_path))
            
            # 8. 動画を結合して最終出力
            final_video_name = f"{self.project_name}.mp4"
            final_video_path = self.base_output_dir / final_video_name
            self.combine_videos(video_paths, str(final_video_path))
            
            # 9. PPTXで使用したPDFファイルの情報をログ出力
            if corresponding_pdf_path:
                logger.info(f"Used corresponding PDF: {corresponding_pdf_path}")
            
            logger.info(f"Presentation video generation completed: {final_video_path}")
            return str(final_video_path)
            
        except Exception as e:
            logger.error(f"Error in presentation video generation: {e}")
            raise


def main():
    """メイン処理"""
    if len(sys.argv) < 2:
        print("使用方法: python main.py <pptxファイルパス>")
        print("対応形式: PowerPoint (.pptx) のみ")
        sys.exit(1)
    
    file_path = sys.argv[1]
    
    if not os.path.exists(file_path):
        print(f"エラー: ファイルが見つかりません: {file_path}")
        sys.exit(1)
    
    # ファイル形式をチェック（PPTX のみ）
    file_extension = os.path.splitext(file_path)[1].lower()
    if file_extension != '.pptx':
        print(f"エラー: サポートされていないファイル形式です: {file_extension}")
        print("対応形式: PowerPoint (.pptx) のみ")
        sys.exit(1)
    
    # プロジェクト名をファイル名から生成
    project_name = Path(file_path).stem
    
    # 必要なAPI KEYの確認
    if not os.getenv("GOOGLE_API_KEY"):
        print("エラー: 環境変数にGOOGLE_API_KEYが見つかりません")
        sys.exit(1)
    
    generator = PresentationVideoGenerator(project_name)
    
    try:
        final_video = generator.generate_presentation_video(file_path)
        print(f"✅ プレゼンテーション動画が正常に生成されました: {final_video}")
        
        # 処理結果の追加情報
        base_name = Path(file_path).stem
        used_pdf_path = f"slides/pdf/{base_name}.pdf"
        print("📝 補足: 発表者ノートがあるスライドは音声付き、ノートがないスライドは3秒間の無音表示になります")
        print(f"📄 使用したPDFファイル: {used_pdf_path}")
    except Exception as e:
        print(f"❌ エラー: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()